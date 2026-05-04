import streamlit as st
import openpyxl
import pdfplumber
import re
import email
import io
from email import policy
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

# ==========================================
# 0. CONFIGURAÇÃO DA INTERFACE (STREAMLIT)
# ==========================================
st.set_page_config(page_title="Gerador de Handover", page_icon="📄", layout="centered")

# CSS Customizado para deixar a interface com a identidade laranja da empresa
st.markdown("""
    <style>
    /* Cor dos botões principais */
    .stButton>button {
        background-color: #F26522; /* Laranja corporativo */
        color: white;
        border-radius: 5px;
        border: none;
        padding: 10px 24px;
        font-weight: bold;
    }
    .stButton>button:hover {
        background-color: #D9531E; /* Laranja mais escuro no hover */
        color: white;
    }
    /* Cor dos títulos */
    h1, h2, h3 {
        color: #F26522;
    }
    /* Cor da barra de progresso / spinners */
    .stSpinner > div > div {
        border-top-color: #F26522 !important;
    }
    </style>
""", unsafe_allow_html=True)

st.title("Gerador de Handover Automático")
st.markdown("Faça o upload dos arquivos do cliente abaixo. O template oficial será preenchido automaticamente.")

# ==========================================
# 1. FUNÇÕES DE AJUDA
# ==========================================
MSG_PADRAO = "###informação não encontrada nos documentos###"
MSG_MANUAL = "[PREENCHER MANUALMENTE]"

def formatar_moeda(valor):
    if isinstance(valor, (int, float)):
        return f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return str(valor)

def formatar_porcentagem(valor):
    if isinstance(valor, (int, float)):
        return f"{valor * 100:.1f}%".replace(".", ",")
    return str(valor)

def extrair_texto_entre_ancoras(texto, palavra_inicial, palavra_final):
    padrao = f"{re.escape(palavra_inicial)}(.*?){re.escape(palavra_final)}"
    resultado = re.search(padrao, texto, re.DOTALL | re.IGNORECASE)
    if resultado:
        return re.sub(r'\s+', ' ', resultado.group(1)).strip()
    return ""

def resumir_texto(texto, limite=250):
    if not texto or texto == MSG_PADRAO or texto == "": return MSG_PADRAO
    return (texto[:limite] + "...") if len(texto) > limite else texto

def buscar_na_coluna(aba, rotulo, col_busca='C', col_retorno='D'):
    for row in range(1, 50):
        val = aba[f'{col_busca}{row}'].value
        if val and rotulo.lower() in str(val).lower():
            retorno = aba[f'{col_retorno}{row}'].value
            return retorno if retorno else ""
    return ""

def ler_texto_email(file_obj):
    """Extrai o texto puro de um arquivo .eml diretamente do upload"""
    try:
        msg = email.message_from_binary_file(file_obj, policy=policy.default)
        corpo = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    corpo = part.get_content()
                    break
        else:
            corpo = msg.get_content()
        return corpo
    except:
        return ""

# ==========================================
# 2. INTERFACE DE UPLOADS
# ==========================================
col1, col2 = st.columns(2)
with col1:
    arquivo_ficha = st.file_uploader("1. Ficha OI (Excel)", type=["xlsx"])
    arquivo_pdf = st.file_uploader("3. Documento RA (PDF)", type=["pdf"])
with col2:
    arquivo_orcamento = st.file_uploader("2. Orçamento (Excel)", type=["xlsx"])
    arquivo_email = st.file_uploader("4. Plano de Trabalho (E-mail .eml)", type=["eml"])

arquivo_template_ppt = 'Handover-teste.pptx' # Fixo no código

# ==========================================
# 3. MOTOR DE PROCESSAMENTO (EXECUTADO AO CLICAR NO BOTÃO)
# ==========================================
if st.button("Gerar Apresentação"):
    if not (arquivo_ficha and arquivo_orcamento and arquivo_pdf and arquivo_email):
        st.warning("⚠️ Por favor, faça o upload de todos os 4 arquivos para gerar a apresentação.")
    else:
        with st.spinner("Extraindo dados e montando a apresentação..."):
            
            dados_extraidos = {
                "{{Registros}}": MSG_MANUAL, "{{Go-NoGo}}": MSG_MANUAL, "{{ACI}}": MSG_MANUAL,
                "{{C / O / P / O}}": MSG_MANUAL, "{{PC / EAP}}": MSG_MANUAL, "{{OC}}": MSG_MANUAL, "{{CLO}}": MSG_MANUAL,
                "{{ORA}}": MSG_MANUAL, "{{LGR}}": MSG_MANUAL, "{{Bid Margin}}": MSG_MANUAL,
                "{{Contrato em vigência (Contrato nº)}}": MSG_MANUAL, "{{Valor/OnePager}}": MSG_MANUAL,
                "{{Escopo}}": MSG_MANUAL, "{{Responsabilidades e Penalidades}}": MSG_MANUAL,
                "{{Faturamento}}": MSG_MANUAL, "{{Preço / Reajuste}}": MSG_MANUAL, "{{Sustentabilidade}}": MSG_MANUAL
            }

            # --- EXTRAÇÃO EXCEL ---
            nome_gerente = "Gerente de Projetos"
            nome_coord = "Coordenador"

            try:
                f_wb = openpyxl.load_workbook(arquivo_ficha, data_only=True)['Nova Ficha - OI']
                nome_gerente = buscar_na_coluna(f_wb, 'Nome Gerente') or "Gerente de Projetos"
                nome_coord = buscar_na_coluna(f_wb, 'Nome Coordenador') or "Coordenador"
                
                dados_extraidos.update({
                    "{{Cliente}}": f_wb['D22'].value, "{{OSC ID}}": f_wb['D4'].value,
                    "{{Projeto}}": f_wb['D5'].value, "{{Localização}}": f_wb['D19'].value,
                    "{{C. custo}}": f_wb['D48'].value
                })
                
                o_wb = openpyxl.load_workbook(arquivo_orcamento, data_only=True)['INFO']
                v_bruto = o_wb['C11'].value
                r_bruta = o_wb['C19'].value
                imp = (v_bruto - r_bruta) if (v_bruto and r_bruta) else 0
                dro_dias = o_wb['C8'].value
                
                dados_extraidos.update({
                    "{{Val/OP}}": formatar_moeda(v_bruto),
                    "{{Impostos}}": formatar_moeda(imp), "{{IMP}}": formatar_moeda(imp),
                    "{{Terceiros}}": formatar_moeda(o_wb['C20'].value), "{{TRC}}": formatar_moeda(o_wb['C20'].value),
                    "{{Despesas}}": formatar_moeda(o_wb['C21'].value), "{{DPS}}": formatar_moeda(o_wb['C21'].value),
                    "{{Net Revenue}}": formatar_moeda(o_wb['C22'].value), "{{NT}}": formatar_moeda(o_wb['C22'].value),
                    "{{Margem Bruta}}": formatar_moeda(o_wb['C25'].value), "{{MB}}": formatar_porcentagem(o_wb['C28'].value),
                    "{{Pessoas & CEs}}": o_wb['C15'].value, "{{PEC}}": o_wb['C15'].value,
                    "{{DRO (dias)}}": str(dro_dias) if dro_dias else "N/A"
                })
            except Exception as e: st.error(f"Aviso Excel: {e}")

            # --- EXTRAÇÃO PDF ---
            try:
                with pdfplumber.open(arquivo_pdf) as pdf:
                    txt_pdf = "\n".join([p.extract_text() for p in pdf.pages[2:] if p.extract_text()])
                
                dados_extraidos.update({
                    "{{Objetivos do cliente para o projeto}}": resumir_texto(extrair_texto_entre_ancoras(txt_pdf, "1. OBJETIVO", "2. DOCUMENTOS DE REFERÊNCIA"), 200),
                    "{{Premissas de custeio}}": resumir_texto(extrair_texto_entre_ancoras(txt_pdf, "4. PREMISSAS", "5. FORA DE ESCOPO"), 200),
                    "{{Fora de escopo contratual}}": resumir_texto(extrair_texto_entre_ancoras(txt_pdf, "5. FORA DE ESCOPO", "6. PRAZO"), 200)
                })
            except Exception as e: st.error(f"Aviso PDF: {e}")

            # --- EXTRAÇÃO E-MAIL ---
            texto_email = ler_texto_email(arquivo_email)
            if texto_email:
                digi_extraido = extrair_texto_entre_ancoras(texto_email, "Digital:", "Outro Topico:")
                dados_extraidos["{{Digital}}"] = resumir_texto(digi_extraido)
            else:
                dados_extraidos["{{Digital}}"] = MSG_PADRAO

            # --- MOTOR CONTEXT-AWARE ---
            texto_contexto = (str(dados_extraidos.get("{{Escopo}}", "")) + " " + str(dados_extraidos.get("{{Premissas de custeio}}", ""))).lower()

            if "remota" in texto_contexto or "remoto" in texto_contexto or "escritório" in texto_contexto:
                dados_extraidos["{{SSO}}"] = "Atividades 100% remotas (Office/Home). Aplicação das políticas corporativas de ergonomia e saúde mental."
            else:
                dados_extraidos["{{SSO}}"] = "Atividades em campo. Uso obrigatório de EPIs completos, integração no site e cumprimento das Regras de Ouro."

            if "sap" in texto_contexto or "cadastro" in texto_contexto or "ch master" in texto_contexto:
                dados_extraidos["{{Entregáveis [Deliverables]}}"] = "Listas Técnicas, Cadastros em Sistema e Relatórios"
                dados_extraidos["{{Entregáveis}}"] = "Cadastro de Materiais (SAP/CH Master)"
            else:
                dados_extraidos["{{Entregáveis [Deliverables]}}"] = "Relatórios Técnicos de Engenharia e Medições"
                dados_extraidos["{{Entregáveis}}"] = "Relatório de Medição Mensal"

            mb_str = str(dados_extraidos.get("{{MB}}", "0")).replace('%', '').replace(',', '.')
            try: margem = float(mb_str)
            except: margem = 0.0

            if margem >= 30.0:
                dados_extraidos["{{Estratégia de execução}}"] = "Foco em mobilização ágil da equipe técnica e precisão nas entregas para assegurar a alta margem projetada."
                dados_extraidos["{{Riscos e Oportunidades}}"] = "Risco de negócio baixo. Oportunidade de otimização de horas da equipe para aumentar rentabilidade."
            else:
                dados_extraidos["{{Estratégia de execução}}"] = "Gestão severa de horas trabalhadas e controle diário de produtividade da equipe."
                dados_extraidos["{{Riscos e Oportunidades}}"] = "Risco de estouro de horas e compressão de margem. Foco total em evitar retrabalho técnico."

            dados_extraidos["{{Fatores críticos de sucesso}}"] = "Comunicação fluida com o cliente, assertividade técnica nos cadastros/relatórios e cumprimento rigoroso dos prazos."

            if dados_extraidos["{{Digital}}"] == MSG_PADRAO:
                dados_extraidos["{{Digital}}"] = "Uso do ecossistema de dados da Arcadis (SharePoint/Teams) para rastreabilidade e segurança da informação."

            dados_extraidos["{{Recursos Chave [internos]}}"] = "Equipe de Gestão e Projetos"
            dados_extraidos["{{recurso/internos}}"] = f"Gerente: {nome_gerente}\nCoord.: {nome_coord}"
            dados_extraidos["{{comentários/internos}}"] = "Liderança técnica, gestão de recursos do projeto e principal ponto focal de alinhamento com a VALE."

            dados_extraidos.update({
                "{{Data/Prazo/entregáveis}}": "Conforme Cronograma",
                "{{Eventos de faturamento}}": "Medição de Serviços (T&M)",
                "{{Data/Prazo/faturamento}}": "Mensal",
                "{{Valor}}": dados_extraidos.get("{{Val/OP}}", "")
            })

            # --- PROCESSAMENTO DO POWERPOINT ---
            try:
                prs = Presentation(arquivo_template_ppt)
                
                def processar_formas(shapes):
                    for item in shapes:
                        if getattr(item, "shape_type", None) == 6:
                            processar_formas(item.shapes)
                        if getattr(item, "has_table", False):
                            for row in item.table.rows:
                                for cell in row.cells:
                                    processar_formas([cell])
                        
                        if hasattr(item, "text_frame") and item.text_frame:
                            item.text_frame.word_wrap = True
                            item.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            
                            for paragraph in item.text_frame.paragraphs:
                                for chave, valor in dados_extraidos.items():
                                    if valor and chave in paragraph.text:
                                        paragraph.text = paragraph.text.replace(chave, str(valor))
                                
                                if "{{" in paragraph.text:
                                    paragraph.text = re.sub(r'\{\{.*?\}\}', MSG_PADRAO, paragraph.text)
                                
                                paragraph.line_spacing = 1.0
                                paragraph.space_before = Pt(0)
                                paragraph.space_after = Pt(0)
                                
                                for run in paragraph.runs:
                                    run.font.size = Pt(12)

                for slide in prs.slides:
                    processar_formas(slide.shapes)
                    if hasattr(slide, 'slide_layout'): processar_formas(slide.slide_layout.shapes)
                    if hasattr(slide, 'slide_master'): processar_formas(slide.master.shapes)

                # Salvar em memória para o usuário baixar
                ppt_io = io.BytesIO()
                prs.save(ppt_io)
                ppt_io.seek(0)
                
                st.success("✅ Apresentação gerada com sucesso!")
                
                st.download_button(
                    label="📥 Baixar Apresentação (PPTX)",
                    data=ppt_io,
                    file_name="Apresentacao_Automatica_Sossego.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            except FileNotFoundError:
                st.error(f"❌ Erro: O arquivo '{arquivo_template_ppt}' não foi encontrado na mesma pasta do código.")